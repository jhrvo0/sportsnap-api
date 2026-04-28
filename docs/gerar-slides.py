#!/usr/bin/env python3
"""
Gerador de apresentação PPTX para o projeto SportSnap — 1a Entrega.
Instale: pip3 install python-pptx
Execute:  python3 gerar-slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Cores ────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(26, 26, 46)      # #1a1a2e
GREEN     = RGBColor(0, 255, 127)      # #00FF7F
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(136, 136, 136)
DARK_CARD = RGBColor(30, 30, 56)       # card bg
BORDER_G  = RGBColor(0, 200, 100)      # softer green border
LIGHT_BG  = RGBColor(38, 38, 68)       # slightly lighter bg for cards

# ── Dimensoes 16:9 ──────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Blank layout
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs,
                     alignment=PP_ALIGN.LEFT, spacing=None):
    """runs = [(text, size, color, bold, font_name), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if spacing:
        p.space_after = Pt(spacing)
    for i, (txt, sz, col, bld, fn) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = txt
        else:
            run = p.add_run()
            run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = fn or "Calibri"
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines,
                          alignment=PP_ALIGN.LEFT):
    """lines = [(text, size, color, bold, font_name), ...]  — one per paragraph"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, col, bld, fn) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = col
        p.font.bold = bld
        p.font.name = fn or "Calibri"
        p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None,
             border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    return shape

def add_pill(slide, left, top, text, fill=None, border=GREEN,
             text_color=GREEN, font_size=10):
    w, h = Inches(len(text) * 0.09 + 0.4), Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    return shape, w

def add_header_bar(slide, section_label=""):
    # top bar
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.7), fill_color=RGBColor(20, 20, 38))
    # logo
    add_rich_textbox(slide, Inches(0.5), Inches(0.12), Inches(3), Inches(0.5),
                     [("SPORT", 20, WHITE, True, "Calibri"),
                      ("SNAP", 20, GREEN, True, "Calibri")])
    # section
    if section_label:
        add_textbox(slide, Inches(9), Inches(0.15), Inches(4), Inches(0.4),
                    section_label, 14, GRAY, alignment=PP_ALIGN.RIGHT)

def add_footer_bar(slide, left_text="", right_text=""):
    bar = add_rect(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5),
                   fill_color=RGBColor(20, 20, 38))
    if left_text:
        add_textbox(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.4),
                    left_text, 10, GRAY)
    if right_text:
        add_textbox(slide, Inches(9), Inches(7.05), Inches(4), Inches(0.4),
                    right_text, 10, GRAY, alignment=PP_ALIGN.RIGHT)


def add_step_card(slide, left, top, width, step_num, title, desc, fields,
                  accent=GREEN):
    """Draw a card with colored top border, title, description, and fields table."""
    card_h = Inches(0.3 + 0.25 + 0.3 + len(fields) * 0.28 + 0.1)
    # card bg
    add_rect(slide, left, top, width, card_h, fill_color=DARK_CARD)
    # top accent
    add_rect(slide, left, top, width, Inches(0.06), fill_color=accent)
    # step number
    add_textbox(slide, left + Inches(0.2), top + Inches(0.12), Inches(1), Inches(0.3),
                f"PASSO {step_num:02d}", 11, accent, bold=True)
    # title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.38), width - Inches(0.4), Inches(0.3),
                title, 15, WHITE, bold=True)
    # fields
    y = top + Inches(0.72)
    for fname, fval in fields:
        add_textbox(slide, left + Inches(0.2), y, Inches(1.8), Inches(0.25),
                    fname, 10, GREEN, bold=True)
        add_textbox(slide, left + Inches(2.0), y, width - Inches(2.4), Inches(0.25),
                    fval, 10, GRAY)
        y += Inches(0.25)
    return card_h


def add_feature_card(slide, left, top, width, fnum, title, desc):
    h = Inches(1.6)
    add_rect(slide, left, top, width, h, fill_color=DARK_CARD)
    add_rect(slide, left, top, width, Inches(0.06), fill_color=GREEN)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.3),
                f"F{fnum:02d}", 24, GREEN, bold=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), Inches(0.3),
                title, 14, WHITE, bold=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.9), width - Inches(0.3), Inches(0.6),
                desc, 10, GRAY)
    return h


def journey_title_slide(slide, num, title, quote, labels, persona="Atleta",
                         section="", slide_num=0):
    set_slide_bg(slide)
    add_header_bar(slide, section)
    add_footer_bar(slide, f"JORNADA {num:02d}", f"{slide_num:02d} — JORNADA {num:02d}")
    # Big number
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(2),
                f"{num:02d}", 120, GREEN, bold=True)
    # Title
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(7), Inches(1),
                title, 32, WHITE, bold=True)
    # Quote
    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(8), Inches(1),
                f'"{quote}"', 16, GRAY)
    # Labels
    x = Inches(0.8)
    for lbl in labels:
        pill, pw = add_pill(slide, x, Inches(5.5), lbl)
        x += pw + Inches(0.15)
    # Persona
    add_textbox(slide, Inches(9.5), Inches(3.5), Inches(3), Inches(0.3),
                "PERSONA", 10, GRAY, bold=True)
    add_textbox(slide, Inches(9.5), Inches(3.8), Inches(3), Inches(0.5),
                persona, 22, WHITE, bold=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)

# Title
add_rich_textbox(s, Inches(0.8), Inches(1.5), Inches(10), Inches(1.5),
                 [("Sport", 60, WHITE, True, "Calibri"),
                  ("Snap", 60, GREEN, True, "Calibri")])

# Subtitle
add_textbox(s, Inches(0.8), Inches(3.0), Inches(10), Inches(0.8),
            "Performance esportiva real + Fotografia profissional\n"
            "em um unico ecossistema.",
            22, GRAY)

# Tags
tags = ["ENGENHARIA DE REQUISITOS", "DOMAIN-DRIVEN DESIGN", "BDD",
        "CUCUMBER + JUNIT 5"]
x = Inches(0.8)
for t in tags:
    pill, pw = add_pill(s, x, Inches(4.5), t)
    x += pw + Inches(0.15)

# Team
add_textbox(s, Inches(0.8), Inches(5.5), Inches(6), Inches(0.3),
            "Antonio Paes  |  Galileu Moares  |  Marco Maciel  |  jhrvo0",
            12, GRAY)
add_textbox(s, Inches(0.8), Inches(5.85), Inches(6), Inches(0.3),
            "CESAR School  —  Eng. de Requisitos + CCPD  —  2026",
            11, GRAY)

# Footer
add_footer_bar(s, "PRESSIONE  \u2192  PARA COMECAR", "01 — CAPA")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O DOMINIO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "DOMINIO")
add_footer_bar(s, "LINGUAGEM UBIQUA", "02 — DOMINIO")

add_textbox(s, Inches(0.8), Inches(1.2), Inches(7), Inches(1),
            "Atletas treinam no escuro e\ndesconectados de suas conquistas.",
            30, WHITE, bold=True)

add_textbox(s, Inches(0.8), Inches(2.6), Inches(7), Inches(1.5),
            "O ecossistema SportSnap une performance esportiva real com fotografia\n"
            "profissional. Atletas acumulam XP oculto (Shadow Stats) durante treinos\n"
            "que so se torna oficial quando adquirem uma licenca de imagem — o Reveal.\n"
            "Isso cria um ciclo de engajamento entre treinar, comprar e evoluir.",
            14, GRAY)

# Metrics row
metrics = [("5", "JORNADAS\nMAPEADAS"), ("8+", "FUNCIONA-\nLIDADES"),
           ("3", "SUB-\nDOMINIOS"), ("\u221e", "POSSIBI-\nLIDADES")]
x = Inches(0.8)
for big, label in metrics:
    add_textbox(s, x, Inches(4.4), Inches(1.8), Inches(0.8),
                big, 48, GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(5.3), Inches(1.8), Inches(0.6),
                label, 10, GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(2.2)

# Subdomains (right)
add_textbox(s, Inches(9.5), Inches(1.5), Inches(3.5), Inches(0.4),
            "SUBDOMINIOS", 12, GREEN, bold=True)

subdomains = [
    ("Gamificacao", "Core Domain", "Shadow Stats, Reveal, Ranking, Overall"),
    ("Marketplace", "Supporting", "Fotos, Licencas, Split Financeiro"),
    ("Sessao", "Generic", "Spots, Sessoes, Check-ins, Match"),
]
y = Inches(2.1)
for name, typ, desc in subdomains:
    add_rect(s, Inches(9.3), y, Inches(3.7), Inches(1.1), fill_color=DARK_CARD)
    add_rect(s, Inches(9.3), y, Inches(0.06), Inches(1.1), fill_color=GREEN)
    add_rich_textbox(s, Inches(9.6), y + Inches(0.08), Inches(3.3), Inches(0.3),
                     [(name, 14, WHITE, True, "Calibri"),
                      (f"  {typ}", 10, GRAY, False, "Calibri")])
    add_textbox(s, Inches(9.6), y + Inches(0.45), Inches(3.3), Inches(0.5),
                desc, 10, GRAY)
    y += Inches(1.2)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — JORNADA 01 TITULO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
journey_title_slide(s, 1,
    "Treinar e marcar presenca.",
    "Quero fazer check-in no local e registrar meu treino para acumular progresso.",
    ["JORNADA", "SESSAO & ATIVIDADE"],
    persona="Atleta",
    section="JORNADA 01", slide_num=3)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — JORNADA 01 DECOMPOSICAO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 01 — DECOMPOSICAO")
add_footer_bar(s, "VALIDARCHECKINIMPL  ·  REGISTRODEATIVIDADE", "04 — JORNADA 01")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "2 passos com regras rigidas de dominio.", 28, WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
            "Cada passo valida campos obrigatorios e aplica invariantes do contexto de Sessao.",
            14, GRAY)

add_step_card(s, Inches(0.8), Inches(2.5), Inches(5.5), 1,
              "Check-in Georreferenciado",
              "",
              [("atletaId", "obrigatorio — UUID do atleta"),
               ("sessionId", "obrigatorio — sessao ativa no Spot"),
               ("latitude", "validacao de range (-90 a 90)"),
               ("longitude", "validacao de range (-180 a 180)"),
               ("timestampCheckIn", "automatico — Instant.now()")])

add_step_card(s, Inches(7.0), Inches(2.5), Inches(5.5), 2,
              "Registro de Atividade",
              "",
              [("distancia", "double — km percorridos"),
               ("duracao", "Duration — tempo total"),
               ("intensidade", "enum — LEVE / MODERADA / INTENSA"),
               ("tipoEsporte", "enum — CORRIDA / MUSCULACAO / ...")])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — JORNADA 01 FUNCIONALIDADES (F01 + F02)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 01 — FUNCIONALIDADES")
add_footer_bar(s, "F01 + F02  ·  NAO-TRIVIAIS", "05 — JORNADA 01")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "Funcionalidades da Jornada 01", 28, WHITE, bold=True)

add_feature_card(s, Inches(0.8), Inches(2.0), Inches(5.5), 1,
                 "CHECK-IN GEORREFERENCIADO",
                 "Validacao de localizacao via GPS + janela temporal.\n"
                 "Garante que o atleta esta fisicamente no Spot\n"
                 "dentro do horario da sessao ativa.")

add_feature_card(s, Inches(7.0), Inches(2.0), Inches(5.5), 2,
                 "CALCULO DE SHADOW STATS",
                 "Acumulo oculto de XP baseado no registro de atividade.\n"
                 "O XP fica em StatusPotencial ate o atleta\n"
                 "disparar o Reveal com uma licenca de imagem.")

# Wireframe placeholder
add_rect(s, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5),
         fill_color=LIGHT_BG, border_color=GRAY)
add_textbox(s, Inches(1.5), Inches(4.8), Inches(4), Inches(0.5),
            "[  PROTOTIPO F01  ]", 14, GRAY, alignment=PP_ALIGN.CENTER)

add_rect(s, Inches(7.0), Inches(4.0), Inches(5.5), Inches(2.5),
         fill_color=LIGHT_BG, border_color=GRAY)
add_textbox(s, Inches(7.7), Inches(4.8), Inches(4), Inches(0.5),
            "[  PROTOTIPO F02  ]", 14, GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — JORNADA 02 TITULO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
journey_title_slide(s, 2,
    "Descobrindo suas fotos: match automatico.",
    "Quero encontrar fotos minhas automaticamente, cruzando meu check-in com o EXIF das fotos.",
    ["JORNADA", "MATCH & BUSCA"],
    persona="Atleta",
    section="JORNADA 02", slide_num=6)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — JORNADA 02 DECOMPOSICAO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 02 — DECOMPOSICAO")
add_footer_bar(s, "MOTORMATCHIMPL  ·  BUSCAFILTRADAIMPL", "07 — JORNADA 02")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "2 passos: cruzamento automatico + busca refinada.", 28, WHITE, bold=True)

add_step_card(s, Inches(0.8), Inches(2.2), Inches(5.5), 1,
              "Motor de Match Automatico",
              "",
              [("checkInTimestamp", "cruzamento com EXIF datetime"),
               ("spotId", "localizacao do Spot vs foto GPS"),
               ("tolerancia", "janela temporal configuravel"),
               ("score", "grau de confianca do match")])

add_step_card(s, Inches(7.0), Inches(2.2), Inches(5.5), 2,
              "Busca Filtrada Inteligente",
              "",
              [("tipoEsporte", "filtro por modalidade"),
               ("fotografoId", "filtro por fotografo"),
               ("dataRange", "intervalo de datas"),
               ("ordenacao", "relevancia / data / preco")])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — JORNADA 02 FUNCIONALIDADES (F03 + F04)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 02 — FUNCIONALIDADES")
add_footer_bar(s, "F03 + F04  ·  NAO-TRIVIAIS", "08 — JORNADA 02")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "Funcionalidades da Jornada 02", 28, WHITE, bold=True)

add_feature_card(s, Inches(0.8), Inches(2.0), Inches(5.5), 3,
                 "MOTOR DE MATCH AUTOMATICO",
                 "Cruza timestamp do check-in com EXIF da foto.\n"
                 "Calcula score de confianca baseado em proximidade\n"
                 "temporal e geografica.")

add_feature_card(s, Inches(7.0), Inches(2.0), Inches(5.5), 4,
                 "BUSCA FILTRADA INTELIGENTE",
                 "Busca multi-criterio por esporte, fotografo,\n"
                 "data e relevancia. Suporta paginacao e\n"
                 "ordenacao dinamica.")

add_rect(s, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5),
         fill_color=LIGHT_BG, border_color=GRAY)
add_textbox(s, Inches(1.5), Inches(4.8), Inches(4), Inches(0.5),
            "[  PROTOTIPO F03  ]", 14, GRAY, alignment=PP_ALIGN.CENTER)

add_rect(s, Inches(7.0), Inches(4.0), Inches(5.5), Inches(2.5),
         fill_color=LIGHT_BG, border_color=GRAY)
add_textbox(s, Inches(7.7), Inches(4.8), Inches(4), Inches(0.5),
            "[  PROTOTIPO F04  ]", 14, GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — JORNADA 03 TITULO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
journey_title_slide(s, 3,
    "Comprando fotos e revelando evolucao.",
    "Quero comprar minhas fotos e sincronizar minha carta para subir no ranking.",
    ["JORNADA", "MARKETPLACE & GAMIFICACAO"],
    persona="Atleta",
    section="JORNADA 03", slide_num=9)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — JORNADA 03 DECOMPOSICAO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 03 — DECOMPOSICAO")
add_footer_bar(s, "CARRINHOIMPL  ·  SINCRONIZACAOIMPL  ·  OVERALLIMPL", "10 — JORNADA 03")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "3 passos: comprar, sincronizar, calcular.", 28, WHITE, bold=True)

add_step_card(s, Inches(0.4), Inches(2.2), Inches(3.8), 1,
              "Carrinho com Preco Dinamico",
              "",
              [("itens", "lista de FotoId"),
               ("desconto", "pacote automatico"),
               ("precoFinal", "calculado em runtime")])

add_step_card(s, Inches(4.6), Inches(2.2), Inches(3.8), 2,
              "Sincronizacao da Carta (Reveal)",
              "",
              [("atletaId", "obrigatorio"),
               ("licencaId", "validacao pos-compra"),
               ("xpTransferido", "StatusPotencial -> Carta")])

add_step_card(s, Inches(8.8), Inches(2.2), Inches(4.1), 3,
              "Calculo de Overall",
              "",
              [("atributos", "media ponderada"),
               ("pesos", "por tipo de esporte"),
               ("overall", "0.0 a 99.0")])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — JORNADA 03 FUNCIONALIDADES (F05 + F06 + F07)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 03 — FUNCIONALIDADES")
add_footer_bar(s, "F05 + F06 + F07  ·  NAO-TRIVIAIS", "11 — JORNADA 03")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "Funcionalidades da Jornada 03", 28, WHITE, bold=True)

add_feature_card(s, Inches(0.4), Inches(2.0), Inches(3.8), 5,
                 "CARRINHO C/ PRECO DINAMICO",
                 "Desconto automatico por volume.\n"
                 "Regras de pacote aplicadas em\n"
                 "tempo de checkout.")

add_feature_card(s, Inches(4.6), Inches(2.0), Inches(3.8), 6,
                 "SINCRONIZACAO DA CARTA",
                 "Reveal: transfere XP oculto para\n"
                 "a CartaOficial. Requer licenca\n"
                 "valida pos-compra.")

add_feature_card(s, Inches(8.8), Inches(2.0), Inches(4.1), 7,
                 "CALCULO DE OVERALL",
                 "Media ponderada dos atributos\n"
                 "esportivos. Pesos variam por\n"
                 "tipo de esporte.")

add_rect(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.5),
         fill_color=LIGHT_BG, border_color=GRAY)
add_textbox(s, Inches(4), Inches(4.8), Inches(6), Inches(0.5),
            "[  PROTOTIPOS F05 / F06 / F07  ]", 14, GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — JORNADA 04 TITULO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
journey_title_slide(s, 4,
    "Monetizando: o lado do fotografo.",
    "Quero fazer upload das minhas fotos em lote e ter elas indexadas automaticamente.",
    ["JORNADA", "OPERACAO DO FOTOGRAFO"],
    persona="Fotografo",
    section="JORNADA 04", slide_num=12)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — JORNADA 04 DECOMPOSICAO + F08
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 04 — DECOMPOSICAO & FUNCIONALIDADE")
add_footer_bar(s, "UPLOADLOTEIMPL  ·  EXIFEXTRACTIONSERVICE", "13 — JORNADA 04")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "1 passo com processamento paralelo.", 28, WHITE, bold=True)

add_step_card(s, Inches(0.8), Inches(2.2), Inches(5.5), 1,
              "Upload e Indexacao em Lote",
              "",
              [("fotografoId", "obrigatorio — UUID"),
               ("arquivos[]", "multipart — batch upload"),
               ("exifExtraction", "paralelo — CompletableFuture"),
               ("spotDetection", "GPS do EXIF -> Spot mais proximo"),
               ("indexacao", "elastica — tags automaticas")])

add_feature_card(s, Inches(7.0), Inches(2.2), Inches(5.5), 8,
                 "UPLOAD E INDEXACAO EM LOTE",
                 "Extracao paralela de metadados EXIF.\n"
                 "Associacao automatica de fotos a Spots\n"
                 "via coordenadas GPS das fotos.")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — JORNADA 05 TITULO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
journey_title_slide(s, 5,
    "Competindo: ranking e engajamento.",
    "Quero ver minha posicao no ranking local e competir com outros atletas do meu Spot.",
    ["JORNADA", "ENGAJAMENTO SOCIAL"],
    persona="Atleta",
    section="JORNADA 05", slide_num=14)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — JORNADA 05 DECOMPOSICAO + F09
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "JORNADA 05 — DECOMPOSICAO & FUNCIONALIDADE")
add_footer_bar(s, "RANKINGPROXY  ·  RANKINGITERATOR", "15 — JORNADA 05")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "1 passo: ranking dinamico local.", 28, WHITE, bold=True)

add_step_card(s, Inches(0.8), Inches(2.2), Inches(5.5), 1,
              "Ranking Dinamico Local",
              "",
              [("spotId", "escopo local por Spot"),
               ("tipoEsporte", "filtro por modalidade"),
               ("sincronizados", "apenas cartas sincronizadas"),
               ("ordenacao", "overall DESC"),
               ("cache", "Proxy pattern — TTL configuravel")])

add_feature_card(s, Inches(7.0), Inches(2.2), Inches(5.5), 9,
                 "RANKING DINAMICO LOCAL",
                 "Apenas cartas sincronizadas participam.\n"
                 "Usa Iterator para percorrer o ranking\n"
                 "e Proxy para cache com TTL.")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — ARQUITETURA LIMPA
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "ARQUITETURA LIMPA")
add_footer_bar(s, "CLEAN ARCHITECTURE  ·  DDD", "16 — ARQUITETURA")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "Separacao em camadas: dominio no centro.", 28, WHITE, bold=True)

# Layer boxes — concentric-style (stacked cards)
layers = [
    ("DOMAIN (centro)", GREEN, [
        "entities/  —  Atleta, CartaOficial, StatusPotencial, Foto, Spot ...",
        "usecases/  —  SincronizarCartaAtleta, CalcularOverall, ValidarCheckIn ...",
        "repositories/  —  interfaces: AtletaRepository, FotoRepository ..."
    ]),
    ("APPLICATION", RGBColor(0, 200, 100), [
        "implementations/  —  SincronizarCartaAtletaImpl, CalcularOverallImpl ...",
        "strategies/  —  EstrategiaCalculoXpCorrida, EstrategiaCalculoXpMusculacao",
        "templates/  —  TemplateSincronizacao (Template Method pattern)"
    ]),
    ("INFRASTRUCTURE", GRAY, [
        "web/controllers/  —  REST endpoints (Spring Boot)",
        "persistence/  —  JpaAtletaRepository, JpaCartaOficialRepository ...",
        "config/  —  SecurityConfig, CorsConfig, SwaggerConfig"
    ]),
]

y = Inches(1.9)
for title, color, items in layers:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(1.4), fill_color=DARK_CARD,
             border_color=color, border_width=Pt(2))
    add_textbox(s, Inches(1.0), y + Inches(0.05), Inches(4), Inches(0.3),
                title, 14, color, bold=True)
    iy = y + Inches(0.35)
    for item in items:
        add_textbox(s, Inches(1.2), iy, Inches(11), Inches(0.28),
                    item, 11, GRAY)
        iy += Inches(0.3)
    y += Inches(1.55)

# Dependency rule arrow text
add_textbox(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.3),
            "Regra de Dependencia:  Infrastructure  \u2192  Application  \u2192  Domain  (nunca o inverso)",
            12, GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DDD NIVEIS
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "DDD — NIVEIS")
add_footer_bar(s, "DOMAIN-DRIVEN DESIGN", "17 — DDD")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "4 niveis de DDD aplicados ao SportSnap.", 28, WHITE, bold=True)

ddd_levels = [
    ("01", "PRELIMINAR", "Descricao do dominio, linguagem ubiqua definida.\n"
     "Termos: Atleta, Spot, Shadow Stats, Reveal, CartaOficial, Overall, Ranking."),
    ("02", "ESTRATEGICO", "3 Bounded Contexts mapeados:\n"
     "  Core — Gamificacao  |  Supporting — Marketplace  |  Generic — Sessao\n"
     "Context Map com relacoes Upstream/Downstream."),
    ("03", "TATICO", "Entidades: Atleta, CartaOficial, StatusPotencial, Foto, Spot, Sessao\n"
     "Value Objects: AtributoEsportivo, Coordenada, ExifData\n"
     "Aggregates: Atleta (raiz), Sessao (raiz)  |  Repositories: interfaces no domain/"),
    ("04", "OPERACIONAL", "Modelo CML (Context Mapper Language): docs/sportsnap.cml\n"
     "Microservicos: gamification-service, session-service, marketplace-service\n"
     "Docker Compose para orquestracao local."),
]

y = Inches(1.8)
for num, title, desc in ddd_levels:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(1.15), fill_color=DARK_CARD)
    add_textbox(s, Inches(1.0), y + Inches(0.08), Inches(0.6), Inches(0.4),
                num, 28, GREEN, bold=True)
    add_textbox(s, Inches(1.7), y + Inches(0.12), Inches(2.5), Inches(0.3),
                title, 16, WHITE, bold=True)
    add_textbox(s, Inches(1.7), y + Inches(0.42), Inches(10.5), Inches(0.7),
                desc, 10, GRAY)
    y += Inches(1.25)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — TESTES BDD INTRO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "TESTES BDD — INTRO")
add_footer_bar(s, "CUCUMBER + JUNIT 5", "18 — BDD")

add_textbox(s, Inches(0.8), Inches(1.2), Inches(10), Inches(1),
            "Sincronizacao da Carta:\nespecificacao & automacao.", 32, WHITE, bold=True)

add_textbox(s, Inches(0.8), Inches(2.8), Inches(10), Inches(0.8),
            "Testes BDD escritos em portugues usando Cucumber com JUnit 5.\n"
            "Padrao Arrange / Act / Assert via anotacoes @Dado @Quando @Entao.",
            16, GRAY)

# Tech pills
techs = ["BDD EM PORTUGUES", "CUCUMBER", "JUNIT 5", "ARRANGE / ACT / ASSERT",
         "@Dado @Quando @Entao"]
x = Inches(0.8)
for t in techs:
    pill, pw = add_pill(s, x, Inches(4.2), t, text_color=GREEN, border=GREEN, font_size=11)
    x += pw + Inches(0.15)

# Feature files list
add_textbox(s, Inches(0.8), Inches(5.2), Inches(10), Inches(0.3),
            "ARQUIVOS DE ESPECIFICACAO:", 12, GREEN, bold=True)
files = [
    "sincronizacao.feature  —  Sincronizacao da Carta do Atleta (2 cenarios)",
    "checkin.feature  —  Check-in Georreferenciado",
    "venda-foto.feature  —  Venda de Foto no Marketplace",
]
y = Inches(5.6)
for f in files:
    add_textbox(s, Inches(1.0), y, Inches(10), Inches(0.25), f, 11, GRAY)
    y += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — TESTES BDD ESPECIFICACAO (Gherkin)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "BDD — ESPECIFICACAO GHERKIN")
add_footer_bar(s, "sincronizacao.feature", "19 — BDD")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
            "2 cenarios cobrem o caminho feliz e as regras de negocio.", 24, WHITE, bold=True)

gherkin = """# language: pt

Funcionalidade: Sincronizacao da Carta do Atleta

  Cenario: Atleta sincroniza carta com licenca valida
    Dado que o Atleta "Joao" possui um CheckIn registrado hoje
    E possui uma LicencaDeImagem adquirida apos o ultimo Reveal
    Quando o Atleta dispara a Sincronizacao
    Entao os StatusPotencial sao transferidos para a CartaOficial
    E o Overall e recalculado
    E a posicao no Ranking e atualizada

  Cenario: Atleta tenta sincronizar sem licenca valida
    Dado que o Atleta "Maria" possui um CheckIn registrado hoje
    E nao possui uma LicencaDeImagem valida
    Quando o Atleta tenta disparar a Sincronizacao
    Entao a sincronizacao e rejeitada
    E a CartaOficial permanece inalterada"""

add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.0),
         fill_color=RGBColor(18, 18, 35), border_color=GREEN, border_width=Pt(1))

lines = gherkin.strip().split("\n")
y = Inches(1.85)
for line in lines:
    col = GRAY
    bld = False
    if line.strip().startswith("Funcionalidade:"):
        col = GREEN; bld = True
    elif line.strip().startswith("Cenario:"):
        col = WHITE; bld = True
    elif any(line.strip().startswith(k) for k in ["Dado", "Quando", "Entao", "E "]):
        col = RGBColor(180, 180, 220)
    elif line.strip().startswith("#"):
        col = GRAY
    add_textbox(s, Inches(0.9), y, Inches(11.5), Inches(0.22),
                line, 11, col, bold=bld, font_name="Consolas")
    y += Inches(0.27)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — TESTES BDD AUTOMACAO (Java Steps)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_header_bar(s, "BDD — AUTOMACAO JAVA")
add_footer_bar(s, "SincronizacaoSteps.java", "20 — BDD")

add_textbox(s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
            "Steps automatizados com Spring + Cucumber.", 24, WHITE, bold=True)

# Left: Arrange
add_rect(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(5.0),
         fill_color=RGBColor(18, 18, 35), border_color=GREEN, border_width=Pt(1))
add_textbox(s, Inches(0.7), Inches(1.8), Inches(3), Inches(0.3),
            "ARRANGE  (@Dado)", 12, GREEN, bold=True)

arrange_code = [
    '@Dado("que o Atleta {string} possui um CheckIn")',
    "public void atletaPossuiCheckInHoje(String nome) {",
    "  atleta = new Atleta(nome, ...);",
    "  atleta = atletaRepository.save(atleta);",
    "  cartaOficial = new CartaOficial(atleta);",
    "  cartaOficial.setOverall(50.0);",
    "  cartaOficial = cartaOficialRepository.save(...);",
    "  statusPotencial = new StatusPotencial(atleta);",
    "  statusPotencial.setXpAcumulado(100.0);",
    "  statusPotencial.setStreakDeConsistencia(3);",
    "  overallAntes = cartaOficial.getOverall();",
    "}",
]
y = Inches(2.15)
for line in arrange_code:
    col = GREEN if line.startswith("@") else RGBColor(180, 180, 220)
    add_textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.22),
                line, 9, col, font_name="Consolas")
    y += Inches(0.24)

# Right: Act + Assert
add_rect(s, Inches(6.8), Inches(1.7), Inches(6.1), Inches(5.0),
         fill_color=RGBColor(18, 18, 35), border_color=GREEN, border_width=Pt(1))
add_textbox(s, Inches(7.0), Inches(1.8), Inches(5), Inches(0.3),
            "ACT (@Quando)  +  ASSERT (@Entao)", 12, GREEN, bold=True)

assert_code = [
    '@Quando("o Atleta dispara a Sincronizacao")',
    "public void atletaDisparaSincronizacao() {",
    "  sincronizarCartaAtleta.executar(atleta.getId());",
    "}",
    "",
    '@Entao("os StatusPotencial sao transferidos")',
    "public void statusTransferidos() {",
    "  StatusPotencial s = statusPotencialRepo",
    "    .findByAtletaId(atleta.getId()).orElseThrow();",
    '  assertEquals(0.0, s.getXpAcumulado(),',
    '    "XP zerado apos sincronizacao");',
    "}",
    "",
    '@E("o Overall e recalculado")',
    "public void overallRecalculado() {",
    "  CartaOficial c = cartaOficialRepo",
    "    .findByAtletaId(atleta.getId()).orElseThrow();",
    "  assertTrue(c.getOverall() > overallAntes);",
    "}",
]
y = Inches(2.15)
for line in assert_code:
    col = GREEN if line.startswith("@") else RGBColor(180, 180, 220)
    if not line:
        y += Inches(0.12)
        continue
    add_textbox(s, Inches(7.0), y, Inches(5.8), Inches(0.22),
                line, 9, col, font_name="Consolas")
    y += Inches(0.22)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — ENCERRAMENTO
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)

add_rich_textbox(s, Inches(2), Inches(1.5), Inches(9), Inches(1),
                 [("Esporte e evolucao. ", 36, WHITE, False, "Calibri"),
                  ("SportSnap", 36, GREEN, True, "Calibri")],
                 alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(2), Inches(3.0), Inches(9), Inches(1),
            "Obrigado!", 48, GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(2), Inches(4.5), Inches(9), Inches(0.5),
            "Perguntas?", 24, WHITE, alignment=PP_ALIGN.CENTER)

# Team
add_textbox(s, Inches(2), Inches(5.5), Inches(9), Inches(0.3),
            "Antonio Paes  |  Galileu Moares  |  Marco Maciel  |  jhrvo0",
            14, GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(2), Inches(6.0), Inches(9), Inches(0.3),
            "CESAR School  —  Eng. de Requisitos + CCPD  —  2026",
            12, GRAY, alignment=PP_ALIGN.CENTER)

# Footer
add_footer_bar(s, "SPORTSNAP  ·  2026", "21 — FIM")


# ── Salvar ───────────────────────────────────────────────────────────────
import os
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "sportsnap-apresentacao.pptx")
prs.save(output_path)
print(f"Apresentacao salva em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
